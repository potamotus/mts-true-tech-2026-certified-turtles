from __future__ import annotations

import io
import os
import re
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Literal
from urllib.error import URLError
from urllib.request import Request, urlopen

from pptx import Presentation
from pptx.util import Inches, Pt

SlideKind = Literal["content", "section", "thanks", "image"]


@dataclass(frozen=True)
class Slide:
    title: str
    bullets: tuple[str, ...]
    notes: str = ""
    kind: SlideKind = "content"
    image_url: str | None = None


def _storage_dir() -> Path:
    root = os.environ.get("GENERATED_FILES_DIR", "/tmp/certified_turtles_files")
    path = Path(root)
    path.mkdir(parents=True, exist_ok=True)
    return path


_SAFE_FILENAME = re.compile(r"[^A-Za-z0-9._-]+")


def _slugify(text: str, *, fallback: str = "presentation") -> str:
    base = _SAFE_FILENAME.sub("-", text.strip()).strip("-")
    return base or fallback


def _theme_path() -> Path | None:
    env = os.environ.get("PPTX_THEME_PATH")
    if env:
        p = Path(env).expanduser()
        return p if p.is_file() else None
    # packaged next to certified_turtles package: .../certified_turtles/assets/theme0.pptx
    here = Path(__file__).resolve()
    p = here.parent.parent / "assets" / "theme0.pptx"
    return p if p.is_file() else None


def _delete_all_slides(prs: Presentation) -> None:
    """Удаляет все слайды шаблона (как в pptx_gen из AI_bot_with_files)."""
    for i in range(len(prs.slides) - 1, -1, -1):
        r_id = prs.slides._sldIdLst[i].rId  # noqa: SLF001
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[i]  # noqa: SLF001


def _fill_bullet_placeholder(body_ph, bullets: tuple[str, ...], *, font_pt: int = 20) -> None:
    tf = body_ph.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = bullet
        para.font.size = Pt(font_pt)
        para.level = 0


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle or ""


def _add_content_slide(prs: Presentation, slide: Slide) -> None:
    layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.title or ""
    if len(s.placeholders) > 1:
        _fill_bullet_placeholder(s.placeholders[1], slide.bullets)
    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def _add_section_slide(prs: Presentation, slide: Slide) -> None:
    layout = prs.slide_layouts[2]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.title or ""
    if slide.bullets and len(s.placeholders) > 1:
        _fill_bullet_placeholder(s.placeholders[1], slide.bullets, font_pt=18)
    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def _looks_like_raster_image(data: bytes) -> bool:
    if len(data) < 12:
        return False
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return True
    if data[:3] == b"\xff\xd8\xff":
        return True
    if data[:6] in (b"GIF87a", b"GIF89a"):
        return True
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return True
    if data[:2] == b"BM":
        return True
    return False


def _download_image_bytes(url: str, *, max_bytes: int = 4_000_000, timeout: int = 20) -> bytes | None:
    if not url.startswith(("http://", "https://")):
        return None
    from certified_turtles.tools.fetch_url import _is_safe_url

    if not _is_safe_url(url):
        return None
    req = Request(
        url,
        headers={
            "User-Agent": "certified-turtles-pptx/1.0",
            "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
        },
    )
    try:
        with urlopen(req, timeout=timeout) as resp:  # noqa: S310
            ctype = (resp.headers.get("Content-Type") or "").split(";")[0].strip().lower()
            if ctype and (ctype.startswith("text/") or "html" in ctype or ctype == "application/json"):
                return None
            data = resp.read(max_bytes + 1)
    except (OSError, URLError, ValueError):
        return None
    if len(data) > max_bytes:
        return None
    if not _looks_like_raster_image(data):
        return None
    return data


def _add_picture_slide(prs: Presentation, slide: Slide) -> None:
    try:
        layout = prs.slide_layouts[8]
    except IndexError:
        # Theme has no picture layout — fall back to content slide
        _add_content_slide(prs, slide)
        return
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = slide.title or ""
    caption_ph = s.placeholders[2] if len(s.placeholders) > 2 else None
    caption_bullets = slide.bullets if slide.bullets else (slide.title or "Иллюстрация",)
    if caption_ph is not None:
        _fill_bullet_placeholder(caption_ph, caption_bullets, font_pt=16)
    pic_ph = s.placeholders[1]
    raw: bytes | None = None
    if slide.image_url:
        raw = _download_image_bytes(slide.image_url.strip())
    if raw:
        try:
            pic_ph.insert_picture(io.BytesIO(raw))
        except Exception:  # noqa: BLE001
            if caption_ph is not None:
                extra = (
                    f"Не удалось вставить файл по ссылке (проверь прямую ссылку на png/jpg/webp): "
                    f"{slide.image_url}"
                )
                _fill_bullet_placeholder(caption_ph, (*caption_bullets, extra), font_pt=12)
    elif caption_ph is not None and slide.image_url:
        _fill_bullet_placeholder(
            caption_ph,
            (
                *caption_bullets,
                f"(изображение не загрузилось: возможна защита от hotlink или не картинка) — {slide.image_url}",
            ),
            font_pt=12,
        )
    if slide.notes:
        s.notes_slide.notes_text_frame.text = slide.notes


def _build_from_template(title: str, subtitle: str, slides: list[Slide]) -> Presentation:
    path = _theme_path()
    if path is None:
        raise FileNotFoundError("theme0.pptx")
    prs = Presentation(path)
    _delete_all_slides(prs)
    _add_title_slide(prs, title, subtitle)
    for slide in slides:
        k = slide.kind
        if k == "content":
            _add_content_slide(prs, slide)
        elif k in ("section", "thanks"):
            _add_section_slide(prs, slide)
        elif k == "image":
            _add_picture_slide(prs, slide)
        else:
            _add_content_slide(prs, slide)
    return prs


def _build_blank(title: str, subtitle: str, slides: list[Slide]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = subtitle or ""
    bullet_layout = prs.slide_layouts[1]
    for slide in slides:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = slide.title or ""
        body = s.placeholders[1] if len(s.placeholders) > 1 else None
        if body is not None:
            _fill_bullet_placeholder(body, slide.bullets)
        if slide.notes:
            s.notes_slide.notes_text_frame.text = slide.notes
    return prs


def build_presentation(title: str, subtitle: str, slides: list[Slide]) -> Path:
    """Строит .pptx: шаблон theme0.pptx (если есть) или пустая презентация. Возвращает путь к файлу."""
    try:
        prs = _build_from_template(title, subtitle, slides)
    except FileNotFoundError:
        prs = _build_blank(title, subtitle, slides)
    filename = f"{_slugify(title)[:60]}-{uuid.uuid4().hex[:8]}.pptx"
    out = _storage_dir() / filename
    prs.save(out)
    return out
